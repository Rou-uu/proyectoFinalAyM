"""Genera reports/presentacion.pptx con un diseño propio y controlado
(python-pptx): 16:9, tipografía uniforme, cada imagen escalada a su caja sin
deformarse ni salirse, y tablas con estilo. El contenido es la versión NO
técnica de la presentación.

Uso:  python reports/construir_presentacion.py
"""
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

AQUI = Path(__file__).resolve().parent
RAIZ = AQUI.parent
FIG = RAIZ / "figures"
DIAG = RAIZ / "diagrams"

# ---- Paleta (tema "semáforo C5", sobrio) ----
AZUL = RGBColor(0x1F, 0x3A, 0x5F)      # títulos
GRIS = RGBColor(0x33, 0x33, 0x33)      # texto
GRIS_CLARO = RGBColor(0x8A, 0x8A, 0x8A)  # pie
ROJO = RGBColor(0xE0, 0x40, 0x3C)
AMBAR = RGBColor(0xF2, 0xB3, 0x3D)
VERDE = RGBColor(0x3B, 0xA7, 0x57)
AZUL_SUAVE = RGBColor(0xEC, 0xF1, 0xF7)  # fondo de cajas/tabla
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

MX = Inches(0.7)          # margen lateral
CONT_W = SW - 2 * MX      # ancho de contenido


def fondo(slide, color=BLANCO):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def barra_tricolor(slide, y, x=MX, w=CONT_W, h=Inches(0.07)):
    seg = int(w / 3)
    for i, c in enumerate((ROJO, AMBAR, VERDE)):
        r = slide.shapes.add_shape(1, x + i * seg, y, seg, h)
        r.fill.solid(); r.fill.fore_color.rgb = c
        r.line.fill.background()


def _run(parag, texto, size, color, bold, italic):
    if texto == "":
        return
    run = parag.add_run(); run.text = texto
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic
    run.font.name = "Calibri"


_FMT = re.compile(r"\*\*(.+?)\*\*|\*(.+?)\*")


def _set_runs(parag, texto, size, color, bold_default=False):
    """Parsea **negritas** y *cursivas* en runs (sin dejar asteriscos)."""
    pos = 0
    for m in _FMT.finditer(texto):
        if m.start() > pos:
            _run(parag, texto[pos:m.start()], size, color, bold_default, False)
        if m.group(1) is not None:        # **negrita**
            _run(parag, m.group(1), size, color, True, False)
        else:                              # *cursiva*
            _run(parag, m.group(2), size, color, bold_default, True)
        pos = m.end()
    if pos < len(texto):
        _run(parag, texto[pos:], size, color, bold_default, False)


def titulo(slide, texto, n):
    tb = slide.shapes.add_textbox(MX, Inches(0.32), CONT_W, Inches(0.85))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    _set_runs(p, texto, 30, AZUL, bold_default=True)
    barra_tricolor(slide, Inches(1.18))
    # pie
    pf = slide.shapes.add_textbox(MX, SH - Inches(0.42), CONT_W, Inches(0.3))
    par = pf.text_frame.paragraphs[0]
    r = par.add_run(); r.text = "Triage de reportes viales del C5  ·  UNAM"
    r.font.size = Pt(9); r.font.color.rgb = GRIS_CLARO
    nb = slide.shapes.add_textbox(SW - Inches(1.1), SH - Inches(0.42), Inches(0.5), Inches(0.3))
    np_ = nb.text_frame.paragraphs[0]; np_.alignment = PP_ALIGN.RIGHT
    rn = np_.add_run(); rn.text = str(n); rn.font.size = Pt(9); rn.font.color.rgb = GRIS_CLARO


def bullets(slide, items, x=MX, y=Inches(1.45), w=CONT_W, h=Inches(5.2), size=18,
            anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for it in items:
        nivel = 0
        txt = it
        if it.startswith("- "):
            nivel = 1; txt = it[2:]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = nivel
        p.space_after = Pt(10 if nivel == 0 else 4)
        p.line_spacing = 1.05
        marca = "•  " if nivel == 0 else "–  "
        rm = p.add_run(); rm.text = marca
        rm.font.size = Pt(size); rm.font.color.rgb = (AZUL if nivel == 0 else GRIS_CLARO)
        rm.font.bold = True
        _set_runs(p, txt, size if nivel == 0 else size - 2, GRIS)


def imagen_fit(slide, ruta, bx, by, bw, bh, valign="center"):
    """Coloca la imagen escalada para CABER en la caja (bx,by,bw,bh) sin
    deformarse, centrada horizontalmente."""
    iw, ih = Image.open(ruta).size
    ar = iw / ih
    box_ar = bw / bh
    if ar > box_ar:
        w = bw; h = int(bw / ar)
    else:
        h = bh; w = int(bh * ar)
    x = bx + (bw - w) // 2
    y = by if valign == "top" else by + (bh - h) // 2
    slide.shapes.add_picture(str(ruta), x, y, w, h)


def tabla(slide, datos, y=Inches(2.0), col_w=None, size=16):
    """datos: lista de filas (primera = encabezado). Centrada."""
    nfilas = len(datos); ncols = len(datos[0])
    if col_w is None:
        col_w = [CONT_W // ncols] * ncols
    tw = sum(col_w)
    x = (SW - tw) // 2
    th = Inches(0.55) * nfilas
    gf = slide.shapes.add_table(nfilas, ncols, x, y, tw, th)
    t = gf.table
    for j, cw in enumerate(col_w):
        t.columns[j].width = cw
    for i, fila in enumerate(datos):
        for j, val in enumerate(fila):
            c = t.cell(i, j)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = Pt(4); c.margin_bottom = Pt(4)
            c.fill.solid()
            c.fill.fore_color.rgb = AZUL if i == 0 else (AZUL_SUAVE if i % 2 else BLANCO)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            _set_runs(p, val, size, BLANCO if i == 0 else GRIS, bold_default=(i == 0))
    return gf


def portada(slide):
    fondo(slide, AZUL)
    barra_tricolor(slide, Inches(2.35), x=Inches(2.0), w=Inches(9.33), h=Inches(0.12))
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.7), Inches(11.33), Inches(2.2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "¿Es real o falso?"; r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = BLANCO
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(6)
    r2 = p2.add_run()
    r2.text = "Distinguir un accidente real de una falsa alarma — y descubrir dónde ocurren"
    r2.font.size = Pt(20); r2.font.color.rgb = RGBColor(0xCF, 0xDC, 0xEC)
    sb = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(11.33), Inches(1.3))
    sf = sb.text_frame; sf.word_wrap = True
    for txt, sz in [("Proyecto Final · Almacenes y Minería de Datos · Facultad de Ciencias, UNAM", 15),
                    ("José Rubén Alfaro González · Cuenta 320516436", 15),
                    ("Junio de 2026", 13)]:
        par = sf.add_paragraph() if sf.paragraphs[0].runs else sf.paragraphs[0]
        par.alignment = PP_ALIGN.CENTER
        rr = par.add_run(); rr.text = txt; rr.font.size = Pt(sz); rr.font.color.rgb = RGBColor(0xDD, 0xE6, 0xF0)


def cierre(slide):
    fondo(slide, AZUL)
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.33), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "¡Gracias!"; r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = BLANCO
    for txt in ["Repositorio:  github.com/Rou-uu/proyectoFinalAyM",
                "Sitio del proyecto:  rou-uu.github.io/proyectoFinalAyM",
                "¿Preguntas?"]:
        par = tf.add_paragraph(); par.alignment = PP_ALIGN.CENTER; par.space_before = Pt(10)
        rr = par.add_run(); rr.text = txt; rr.font.size = Pt(18); rr.font.color.rgb = RGBColor(0xDD, 0xE6, 0xF0)
    barra_tricolor(slide, Inches(6.6), x=Inches(3.0), w=Inches(7.33), h=Inches(0.1))


# Caja de imagen estándar para slides de 2 columnas
IMG_X, IMG_Y, IMG_W, IMG_H = Inches(7.55), Inches(1.5), Inches(5.1), Inches(5.2)
TXT_W = Inches(6.45)


def dos_columnas(slide, n, tit, items, img):
    titulo(slide, tit, n)
    bullets(slide, items, x=MX, y=Inches(1.55), w=TXT_W, h=Inches(5.1),
            anchor=MSO_ANCHOR.TOP)
    imagen_fit(slide, img, IMG_X, IMG_Y, IMG_W, IMG_H, valign="top")


def nueva():
    s = prs.slides.add_slide(BLANK); fondo(s); return s


# ===================== CONTENIDO =====================
n = 1
portada(nueva())

n += 1; s = nueva(); titulo(s, "El problema, en una frase", n)
bullets(s, [
    "El **C5** recibe cientos de miles de reportes de accidentes viales al año.",
    "**Más de 1 de cada 3** resulta ser **falso o una falsa alarma**.",
    "Mandar una unidad a un reporte que no era real **gasta un recurso escaso** y resta capacidad para las emergencias verdaderas.",
    "**Nuestra meta:** estimar, al entrar el reporte, qué tan probable es que sea **real** — y descubrir **en qué calles** se concentran los accidentes.",
])

n += 1; s = nueva(); titulo(s, "Los datos que usamos", n)
bullets(s, [
    "Reportes **reales y públicos** del C5 de la Ciudad de México (2022–2024).",
    "Cerca de **290 mil** reportes después de la limpieza.",
    "De cada uno sabemos: **por qué canal llegó** (911, cámara, botón de auxilio…), **qué tipo de accidente**, **a qué hora** y **dónde**.",
    "Un reporte que **repite** uno ya atendido lo dejamos fuera: eso es trabajo del sistema de despacho, no de la predicción.",
])

n += 1; dos_columnas(nueva(), n, "Hallazgo 1: el canal dice mucho", [
    "Lo que entra por canales **del propio C5** casi siempre es real: **radio y cámara, ~9 de cada 10**; **botón de auxilio, ~9 de cada 10**.",
    "La **llamada al 911** —la gran mayoría— resulta real solo **6 de cada 10**.",
    "**Idea clave:** confiar casi directo en lo institucional y poner la lupa en las llamadas ciudadanas, donde está la duda.",
], FIG / "01_fuente_confirmacion.png")

n += 1; dos_columnas(nueva(), n, "Hallazgo 2: la hora también importa", [
    "De **madrugada y mañana**, la mayoría son reales (**~7 de cada 10**).",
    "Por la **tarde y noche** baja (**~6 de cada 10**)… y es cuando **más reportes llegan**.",
    "Las horas de mayor volumen son también las de **mayor ruido**.",
], FIG / "01_franja_confirmacion.png")

n += 1; s = nueva(); titulo(s, "¿Qué tan bien predice el modelo?", n)
bullets(s, [
    "Elegimos un **bosque de decisiones** (*Random Forest*: muchos árboles que votan). Así se compara contra lo que pasaría **sin modelo**:",
], y=Inches(1.4), h=Inches(0.9), size=18)
tabla(s, [
    ["", "Sin modelo (“todo es real”)", "Con nuestro modelo"],
    ["Aciertos generales", "~63 de cada 100", "~67 de cada 100"],
    ["Falsas alarmas detectadas", "0  (ninguna)", "~2 de cada 3"],
], y=Inches(2.5), col_w=[Inches(3.4), Inches(4.2), Inches(4.2)], size=16)
bullets(s, [
    "El gran salto es ese **filtro de falsas alarmas que antes no existía**.",
    "Una **red neuronal** empata, pero detecta menos falsos: preferimos el bosque.",
], y=Inches(5.3), h=Inches(1.4), size=16)

n += 1; dos_columnas(nueva(), n, "¿En qué se fija para decidir?", [
    "Lo que **más** le ayuda es **qué tipo de accidente** es y **por qué canal** llegó.",
    "Mucho más que la hora o el lugar exactos.",
    "Tiene sentido: una cámara o un botón suelen ir con eventos confirmados; las llamadas son más variadas.",
], FIG / "02_importancia_permutacion.png")

n += 1; s = nueva(); titulo(s, "Jugamos limpio", n)
bullets(s, [
    "Nos cuidamos de **no darle al modelo información que solo se conoce después** de atender el caso (por ejemplo, cuánto tardó en cerrarse).",
    "Si lo hiciéramos, parecería más listo… pero sería **“hacer trampa” mirando el futuro**: en la vida real esa información **no existe** cuando llega el reporte.",
    "Por eso nuestros resultados son **honestos y reproducibles**.",
])

n += 1; dos_columnas(nueva(), n, "¿Dónde ocurren los accidentes?", [
    "Agrupamos los incidentes por **zonas de alta concentración** (método *DBSCAN*).",
    "Aparecen **puntos calientes** sobre **avenidas reconocibles**: Circuito Interior, Calz. Ignacio Zaragoza, Distribuidor Vial San Antonio…",
    "**Dato curioso:** en las avenidas más saturadas, una **mayor proporción resulta falsa** → conviene **verificarlas más**.",
], FIG / "figura_dbscan_vialidades.png")

n += 1; s = nueva(); titulo(s, "Cómo está hecho por dentro", n)
bullets(s, [
    "Organizamos el programa como una **línea de montaje** (patrón de diseño *Pipeline*): cada etapa hace **una tarea** —limpiar, preparar, predecir— y se encadena.",
    "Ventaja: **ordenado, reutilizable y fácil de revisar**. El diagrama es el **mapa** del código.",
], y=Inches(1.5), h=Inches(1.2), size=15, anchor=MSO_ANCHOR.TOP)
imagen_fit(s, DIAG / "uml_clases.png", Inches(1.0), Inches(2.85), Inches(11.33), Inches(4.35))

n += 1; s = nueva(); titulo(s, "Demostración en vivo", n)
bullets(s, [
    "Guardamos el modelo entrenado y lo **volvemos a usar sin re-entrenar**: le damos un reporte nuevo y dice **qué tan probable es que sea real**.",
], y=Inches(1.4), h=Inches(0.9), size=18)
tabla(s, [
    ["Reporte nuevo", "Probabilidad de ser real"],
    ["Cámara · zona centro · choque con lesionados", "81 %  →  atender"],
    ["Llamada al 911 · periferia · de madrugada", "36 %  →  dudoso"],
    ["Botón de auxilio · choque con lesionados", "90 %  →  atender"],
], y=Inches(2.5), col_w=[Inches(7.2), Inches(4.0)], size=16)
bullets(s, [
    "En vivo: cargamos el modelo y predecimos un caso; y abrimos el **mapa interactivo** de la ciudad.",
], y=Inches(5.6), h=Inches(1.0), size=15)

n += 1; s = nueva(); titulo(s, "En resumen", n)
bullets(s, [
    "**Quién y cómo reporta** dice mucho: lo institucional es muy confiable; la llamada al 911 es donde está la duda.",
    "El modelo aporta algo que antes no existía: **un filtro de falsas alarmas** que atrapa ~2 de cada 3, sin perder los reales.",
    "Sabemos **en qué avenidas** se concentran los accidentes y cuáles necesitan más verificación.",
    "**Límites:** dependemos de cómo el C5 etiqueta los reportes y de los datos 2022–2024; faltaría contexto (clima, tráfico, qué dijo la llamada).",
])

n += 1; cierre(nueva())

out = AQUI / "presentacion.pptx"
prs.save(out)
print("Generado:", out, "|", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
